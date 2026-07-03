"""PowerPoint (.pptx) document processor using python-pptx."""

import asyncio
import logging
from datetime import datetime

logger = logging.getLogger(__name__)


def _load_pptx():
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        raise ImportError(
            "python-pptx is required for PPTX processing. Install with: pip install python-pptx"
        ) from None
    return Presentation


class PPTXProcessor:
    """Extract text and metadata from PowerPoint files."""

    async def extract_text(self, file_path: str) -> str:
        """Extract text from a PPTX file."""
        Presentation = _load_pptx()

        def _extract():
            try:
                prs = Presentation(file_path)
            except Exception as e:
                if "password" in str(e).lower() or "encrypted" in str(e).lower():
                    raise ValueError("Password-protected PowerPoint files are not supported") from e
                raise

            text_parts = []

            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_texts = []
                text_parts.append(f"--- Slide {slide_num} ---")

                for shape in slide.shapes:
                    _extract_shape_text(shape, slide_texts)

                # Speaker notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_texts.append(f"[Notes: {notes}]")

                text_parts.extend(slide_texts)

            return "\n".join(text_parts)

        return await asyncio.to_thread(_extract)

    async def extract_metadata(self, file_path: str) -> dict:
        """Best-effort metadata extraction — returns {} on any failure."""
        Presentation = _load_pptx()

        def _extract():
            try:
                prs = Presentation(file_path)
                props = getattr(prs, "core_properties", None)
                if props is None:
                    return {}
                author = (getattr(props, "author", None) or "").strip() or None
                title = (getattr(props, "title", None) or "").strip() or None
                created = getattr(props, "created", None)
                authored_at = created if isinstance(created, datetime) else None
                return {"author": author, "doc_title": title, "authored_at": authored_at}
            except Exception as e:
                logger.debug("PPTX metadata extraction failed for %s: %s", file_path, e)
                return {}

        return await asyncio.to_thread(_extract)


def _extract_shape_text(shape, texts: list[str]):
    """Recursively extract text from a shape."""
    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                texts.append(text)

    if shape.has_table:
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            row_text = "\t".join(cells)
            if row_text.strip():
                texts.append(row_text)

    # Grouped shapes — recurse into sub-shapes
    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            _extract_shape_text(sub_shape, texts)
